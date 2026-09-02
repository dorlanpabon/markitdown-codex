"""Build distributable plugin, standalone skill and synthetic review fixtures."""

import hashlib
import io
import json
import zipfile
from pathlib import Path

from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

ROOT = Path(__file__).resolve().parents[1]
PLUGIN = ROOT / "plugins" / "markitdown"
DIST = ROOT / "dist"


def build_logo():
    image = Image.new("RGB", (512, 512), "#2563EB")
    draw = ImageDraw.Draw(image)
    draw.rounded_rectangle((106, 62, 358, 450), radius=22, fill="white")
    for y, end in ((138, 304), (184, 290), (230, 252)):
        draw.rounded_rectangle((150, y, end, y + 14), radius=7, fill="#2563EB")
    draw.rounded_rectangle((285, 280, 335, 368), radius=4, fill="#2563EB")
    draw.polygon(((258, 348), (362, 348), (310, 402)), fill="#2563EB")
    assets = PLUGIN / "assets"
    assets.mkdir(exist_ok=True)
    image.save(assets / "logo.png")
    image.resize((128, 128), Image.Resampling.LANCZOS).save(assets / "icon.png")


def archive_bytes(path: Path, entries: dict[str, bytes]):
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for name, data in sorted(entries.items()):
            info = zipfile.ZipInfo(name, date_time=(2026, 9, 2, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = 0o644 << 16
            archive.writestr(info, data)


def review_fixtures() -> dict[str, bytes]:
    pdf = io.BytesIO()
    canvas = Canvas(pdf, invariant=True)
    canvas.drawString(72, 720, "Quarterly report 2026")
    canvas.save()
    docx = io.BytesIO()
    with zipfile.ZipFile(docx, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>',
        )
        archive.writestr(
            "word/document.xml",
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            "<w:body><w:p><w:r><w:t>Project document</w:t></w:r></w:p></w:body></w:document>",
        )
    slides = io.BytesIO()
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[0]).shapes.title.text = "Project slides"
    deck.save(slides)
    sheet = io.BytesIO()
    workbook = Workbook()
    workbook.active.append(["Product", "Quantity"])
    workbook.active.append(["Coffee", 12])
    workbook.save(sheet)
    return {
        "sample.pdf": pdf.getvalue(),
        "sample.docx": docx.getvalue(),
        "sample.pptx": slides.getvalue(),
        "sample.xlsx": sheet.getvalue(),
        "unicode.txt": "Resumen: café, Bogotá y conversión.\n".encode(),
    }


def main():
    manifest = json.loads((PLUGIN / ".codex-plugin/plugin.json").read_text())
    version = manifest["version"]
    build_logo()
    DIST.mkdir(exist_ok=True)
    files = {
        path.relative_to(PLUGIN).as_posix(): path.read_bytes()
        for path in PLUGIN.rglob("*")
        if path.is_file()
        and "__pycache__" not in path.parts
        and path.suffix in {".json", ".md", ".py", ".lock", ".png"}
    }
    files["LICENSE"] = (ROOT / "LICENSE").read_bytes()
    archive_bytes(DIST / f"markitdown-plugin-{version}.zip", files)
    skills = {
        name.removeprefix("skills/"): content
        for name, content in files.items()
        if name.startswith("skills/")
    }
    skills["markitdown/LICENSE"] = files["LICENSE"]
    archive_bytes(DIST / f"markitdown-skills-{version}.zip", skills)
    archive_bytes(DIST / "markitdown-review-fixtures.zip", review_fixtures())
    release_files = [
        DIST / f"markitdown-plugin-{version}.zip",
        DIST / f"markitdown-skills-{version}.zip",
        DIST / "markitdown-review-fixtures.zip",
    ]
    (DIST / "SHA256SUMS").write_text(
        "".join(
            f"{hashlib.sha256(path.read_bytes()).hexdigest()}  {path.name}\n"
            for path in release_files
        ),
        encoding="utf-8",
    )
    print(f"Built {len(release_files)} archives and SHA256SUMS in {DIST}")


if __name__ == "__main__":
    main()

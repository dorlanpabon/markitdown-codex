import asyncio
import json
import os
import socket
import subprocess
import sys
import zipfile
from pathlib import Path

import convert
import pytest
from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

ROOT = Path(__file__).resolve().parents[1]
PLUGIN = ROOT / "plugins" / "markitdown"
SCRIPT = PLUGIN / "skills" / "markitdown" / "scripts" / "convert.py"


@pytest.fixture(autouse=True)
def isolated_root(monkeypatch):
    monkeypatch.delenv("MARKITDOWN_ROOT", raising=False)


@pytest.fixture
def text_file(tmp_path):
    path = tmp_path / "prueba con espacios ñ.txt"
    path.write_text("Resumen: café, Bogotá y conversión.\n", encoding="utf-8")
    return path


def test_text_and_pagination(text_file):
    full = convert.convert_document(str(text_file))
    parts = []
    offset = 0
    while offset is not None:
        page = convert.convert_to_markdown(str(text_file), offset, 7)
        assert page.total_chars == len(full)
        parts.append(page.markdown)
        offset = page.next_offset
    assert "".join(parts) == full
    assert "Bogotá" in full


@pytest.mark.parametrize(
    "value",
    [
        "relative.txt",
        "https://example.com/a.pdf",
        "file:///tmp/a.txt",
        "//server/a.txt",
        "\\\\server\\a.txt",
    ],
)
def test_reject_nonlocal_paths(value):
    with pytest.raises(convert.ConversionError):
        convert.convert_document(value)


def test_root_restriction(text_file, tmp_path, monkeypatch):
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    monkeypatch.setenv("MARKITDOWN_ROOT", str(allowed))
    with pytest.raises(convert.ConversionError, match="outside"):
        convert.convert_document(str(text_file))
    inside = allowed / "ok.txt"
    inside.write_text("inside", encoding="utf-8")
    assert "inside" in convert.convert_document(str(inside))


def test_symlink_cannot_escape_root(tmp_path, text_file, monkeypatch):
    allowed = tmp_path / "allowed"
    allowed.mkdir()
    link = allowed / "escape.txt"
    try:
        link.symlink_to(text_file)
    except OSError:
        pytest.skip("Creating symlinks requires OS permission")
    monkeypatch.setenv("MARKITDOWN_ROOT", str(allowed))
    with pytest.raises(convert.ConversionError, match="outside"):
        convert.convert_document(str(link))


def test_missing_and_unsupported_files(tmp_path):
    with pytest.raises(convert.ConversionError):
        convert.convert_document(str(tmp_path / "missing.txt"))
    path = tmp_path / "data.zip"
    path.write_bytes(b"zip")
    with pytest.raises(convert.ConversionError, match="Unsupported"):
        convert.convert_document(str(path))


def test_input_size_limit(text_file, monkeypatch):
    monkeypatch.setattr(convert, "MAX_INPUT_BYTES", 2)
    with pytest.raises(convert.ConversionError, match="limit"):
        convert.convert_document(str(text_file))


def test_office_expansion_limit(tmp_path, monkeypatch):
    path = tmp_path / "large.docx"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", "x" * 100)
    monkeypatch.setattr(convert, "MAX_ARCHIVE_BYTES", 10)
    with pytest.raises(convert.ConversionError, match="expanded"):
        convert.convert_document(str(path))


def test_empty_document_and_invalid_page(tmp_path, text_file):
    path = tmp_path / "empty.txt"
    path.write_text("", encoding="utf-8")
    assert convert.convert_to_markdown(str(path)).warnings
    for offset, size in [(-1, 1), (0, 0), (0, 100_001), (10_000, 1)]:
        with pytest.raises(convert.ConversionError):
            convert.convert_to_markdown(str(text_file), offset, size)


def test_html_does_not_fetch_remote_images(tmp_path, monkeypatch):
    def no_network(*args, **kwargs):
        raise AssertionError("Conversion attempted network access")

    monkeypatch.setattr(socket.socket, "connect", no_network)
    path = tmp_path / "page.html"
    path.write_text(
        '<h1>Report</h1><img src="https://example.com/image.png"><p>Local only</p>',
        encoding="utf-8",
    )
    result = convert.convert_document(str(path))
    assert "Report" in result and "Local only" in result


def test_csv(tmp_path):
    path = tmp_path / "data.csv"
    path.write_text("name,total\nCoffee,12\n", encoding="utf-8")
    result = convert.convert_document(str(path))
    assert "Coffee" in result and "12" in result


def test_pdf(tmp_path):
    path = tmp_path / "sample.pdf"
    pdf = Canvas(str(path))
    pdf.drawString(72, 720, "Quarterly report 2026")
    pdf.save()
    assert "Quarterly report 2026" in convert.convert_document(str(path))


def test_docx(tmp_path):
    path = tmp_path / "sample.docx"
    with zipfile.ZipFile(path, "w") as archive:
        archive.writestr(
            "[Content_Types].xml",
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>',
        )
        archive.writestr(
            "word/document.xml",
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Project document</w:t></w:r></w:p></w:body></w:document>',
        )
    assert "Project document" in convert.convert_document(str(path))


def test_pptx(tmp_path):
    path = tmp_path / "sample.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[0])
    slide.shapes.title.text = "Project slides"
    deck.save(path)
    assert "Project slides" in convert.convert_document(str(path))


def test_xlsx(tmp_path):
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["Product", "Quantity"])
    workbook.active.append(["Coffee", 12])
    workbook.save(path)
    result = convert.convert_document(str(path))
    assert "Coffee" in result and "Quantity" in result


def test_corrupt_file_has_sanitized_error(tmp_path):
    path = tmp_path / "secret.pdf"
    path.write_bytes(b"private-content-not-a-pdf")
    with pytest.raises(convert.ConversionError) as error:
        convert.convert_document(str(path))
    assert "private-content" not in str(error.value)
    assert str(tmp_path) not in str(error.value)


def test_cli_export_and_overwrite_protection(text_file, tmp_path):
    output = tmp_path / "result.md"
    command = [sys.executable, str(SCRIPT), str(text_file), "--output", str(output)]
    first = subprocess.run(command, capture_output=True, timeout=60)
    assert first.returncode == 0, first.stderr
    original = output.read_bytes()
    second = subprocess.run(command, capture_output=True, timeout=60)
    assert second.returncode != 0
    assert output.read_bytes() == original
    assert "Bogotá" in output.read_text(encoding="utf-8")


def test_actual_manifest_stdio_transport(text_file):
    config = json.loads((PLUGIN / ".mcp.json").read_text())["mcpServers"]["markitdown"]

    async def exercise():
        params = StdioServerParameters(
            command=config["command"],
            args=config["args"],
            cwd=str(PLUGIN),
            env={**os.environ, "MARKITDOWN_ROOT": str(text_file.parent)},
        )
        async with stdio_client(params) as (reader, writer):
            async with ClientSession(reader, writer) as session:
                await session.initialize()
                tools = (await session.list_tools()).tools
                assert len(tools) == 1
                tool = tools[0]
                assert tool.name == "convert_to_markdown"
                assert tool.outputSchema
                assert tool.annotations.readOnlyHint is True
                assert tool.annotations.openWorldHint is False
                assert tool.annotations.destructiveHint is False
                result = await session.call_tool(tool.name, {"path": str(text_file)})
                assert not result.isError
                assert "Bogotá" in result.structuredContent["markdown"]
                invalid = await session.call_tool(
                    tool.name, {"path": "https://example.com/test.pdf"}
                )
                assert invalid.isError

    asyncio.run(asyncio.wait_for(exercise(), timeout=120))
